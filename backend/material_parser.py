"""
课件解析模块
Material Parser Module
功能：解析各种格式的课件文件，提取文本内容和知识点
"""

import os
import logging
from typing import Dict, List, Any, Optional
from pathlib import Path
import asyncio
from datetime import datetime

# 文件处理库
try:
    import pypdf
except ImportError:
    pypdf = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import docx
except ImportError:
    docx = None

from ai_service import get_zhipu_service
from database import SessionLocal
from models import CourseMaterial, MaterialContent

logger = logging.getLogger(__name__)


class MaterialParser:
    """课件解析器"""

    def __init__(self):
        self.ai_service = None

    async def parse_material(self, material_id: str, file_path: str) -> Dict[str, Any]:
        """
        解析课件文件

        Args:
            material_id: 课件ID
            file_path: 文件路径

        Returns:
            解析结果字典
        """
        try:
            # 确定文件类型
            file_extension = Path(file_path).suffix.lower()

            # 根据文件类型选择解析方法
            if file_extension == '.pdf':
                content_text = self._parse_pdf(file_path)
            elif file_extension in ['.ppt', '.pptx']:
                content_text = self._parse_ppt(file_path)
            elif file_extension in ['.doc', '.docx']:
                content_text = self._parse_docx(file_path)
            elif file_extension in ['.txt']:
                content_text = self._parse_txt(file_path)
            else:
                return {
                    "success": False,
                    "error": f"不支持的文件类型：{file_extension}"
                }

            if not content_text:
                return {
                    "success": False,
                    "error": "文件内容为空或解析失败"
                }

            # 使用AI提取知识点和结构化内容
            structured_content = await self._extract_knowledge_points(content_text, file_extension)

            # 保存解析结果到数据库
            await self._save_parsed_content(material_id, structured_content)

            return {
                "success": True,
                "material_id": material_id,
                "content_length": len(content_text),
                "chapters_found": len(structured_content.get("chapters", [])),
                "knowledge_points": len(structured_content.get("knowledge_points", [])),
                "parsed_at": datetime.now().isoformat()
            }

        except Exception as e:
            logger.error(f"课件解析失败：{str(e)}")
            return {
                "success": False,
                "error": str(e)
            }

    def _parse_pdf(self, file_path: str) -> str:
        """解析PDF文件"""
        try:
            if pypdf is None:
                raise ImportError("pypdf库未安装")

            text_content = []
            with open(file_path, 'rb') as file:
                pdf_reader = pypdf.PdfReader(file)
                total_pages = len(pdf_reader.pages)

                logger.info(f"开始解析PDF文件，共{total_pages}页")

                for page_num, page in enumerate(pdf_reader.pages, 1):
                    try:
                        page_text = page.extract_text()
                        if page_text.strip():
                            text_content.append(f"【第{page_num}页】\n{page_text}")

                        # 每处理50页记录一次进度
                        if page_num % 50 == 0:
                            logger.info(f"已处理{page_num}/{total_pages}页")

                    except Exception as e:
                        logger.warning(f"解析第{page_num}页失败：{str(e)}")
                        continue

            full_text = "\n\n".join(text_content)
            logger.info(f"PDF解析完成，提取了{len(full_text)}个字符")
            return full_text

        except Exception as e:
            logger.error(f"PDF解析失败：{str(e)}")
            raise

    def _parse_ppt(self, file_path: str) -> str:
        """解析PPT/PPTX文件"""
        try:
            if Presentation is None:
                raise ImportError("python-pptx库未安装")

            prs = Presentation(file_path)
            text_content = []

            logger.info(f"开始解析PPT文件，共{len(prs.slides)}张幻灯片")

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []

                # 提取文本框中的文本
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())

                # 提取表格中的文本
                for shape in slide.shapes:
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_text = " | ".join([cell.text.strip() for cell in row.cells])
                            slide_text.append(row_text)

                if slide_text:
                    text_content.append(f"【第{slide_num}张幻灯片】\n" + "\n".join(slide_text))

            full_text = "\n\n".join(text_content)
            logger.info(f"PPT解析完成，提取了{len(full_text)}个字符")
            return full_text

        except Exception as e:
            logger.error(f"PPT解析失败：{str(e)}")
            raise

    def _parse_docx(self, file_path: str) -> str:
        """解析Word文档"""
        try:
            if docx is None:
                raise ImportError("python-docx库未安装")

            doc = docx.Document(file_path)
            text_content = []

            logger.info("开始解析Word文档")

            # 提取段落文本
            for para_num, paragraph in enumerate(doc.paragraphs, 1):
                if paragraph.text.strip():
                    text_content.append(paragraph.text.strip())

            # 提取表格文本
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join([cell.text.strip() for cell in row.cells])
                    if row_text.strip():
                        text_content.append(row_text)

            full_text = "\n\n".join(text_content)
            logger.info(f"Word文档解析完成，提取了{len(full_text)}个字符")
            return full_text

        except Exception as e:
            logger.error(f"Word文档解析失败：{str(e)}")
            raise

    def _parse_txt(self, file_path: str) -> str:
        """解析文本文件"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                text = file.read()
            logger.info(f"文本文件解析完成，提取了{len(text)}个字符")
            return text

        except UnicodeDecodeError:
            # 尝试其他编码
            try:
                with open(file_path, 'r', encoding='gbk') as file:
                    text = file.read()
                logger.info(f"文本文件解析完成（GBK编码），提取了{len(text)}个字符")
                return text
            except Exception as e:
                logger.error(f"文本文件解析失败：{str(e)}")
                raise
        except Exception as e:
            logger.error(f"文本文件解析失败：{str(e)}")
            raise

    async def _extract_knowledge_points(self, content_text: str, file_type: str) -> Dict[str, Any]:
        """使用AI提取知识点和结构化内容"""
        try:
            if self.ai_service is None:
                self.ai_service = get_zhipu_service()

            # 限制内容长度避免超出token限制
            truncated_content = content_text[:8000] if len(content_text) > 8000 else content_text

            # 调用AI解析内容
            parse_result = await self.ai_service.parse_material_content(
                content=truncated_content,
                material_type=file_type
            )

            # 如果AI返回了知识点，进行章节划分
            if parse_result.get("knowledge_points"):
                chapters = self._split_into_chapters(content_text, parse_result.get("knowledge_points", []))
            else:
                chapters = []

            return {
                "summary": parse_result.get("summary", ""),
                "knowledge_points": parse_result.get("knowledge_points", []),
                "key_concepts": parse_result.get("key_concepts", []),
                "difficulty_level": parse_result.get("difficulty_level", "intermediate"),
                "estimated_study_time": parse_result.get("estimated_study_time", 60),
                "topics": parse_result.get("topics", []),
                "chapters": chapters,
                "full_text": content_text
            }

        except Exception as e:
            logger.error(f"知识点提取失败：{str(e)}")
            # 返回基础结构
            return {
                "summary": content_text[:500] + "..." if len(content_text) > 500 else content_text,
                "knowledge_points": [],
                "key_concepts": [],
                "difficulty_level": "intermediate",
                "estimated_study_time": 60,
                "topics": [],
                "chapters": [],
                "full_text": content_text
            }

    def _split_into_chapters(self, content: str, knowledge_points: List[str]) -> List[Dict[str, Any]]:
        """将内容分割成章节"""
        chapters = []

        try:
            # 简单的章节分割逻辑
            # 根据常见的章节标题模式进行分割
            import re

            # 查找可能的章节标题
            chapter_patterns = [
                r'^第[一二三四五六七八九十\d]+章[\s\u3000]*.+',
                r'^Chapter\s*\d+[\s\u3000]*.+',
                r'^\d+[\.\s][\s\u3000]*.+',
                r'^[一二三四五六七八九十]+[\s\u3000]*、[\s\u3000]*.+'
            ]

            lines = content.split('\n')
            current_chapter = None
            chapter_content = []

            for line in lines:
                is_chapter_title = False

                for pattern in chapter_patterns:
                    if re.match(pattern, line.strip()):
                        # 保存上一章
                        if current_chapter and chapter_content:
                            chapters.append({
                                "title": current_chapter,
                                "content": "\n".join(chapter_content),
                                "knowledge_points": []
                            })

                        current_chapter = line.strip()
                        chapter_content = []
                        is_chapter_title = True
                        break

                if not is_chapter_title:
                    chapter_content.append(line)

            # 保存最后一章
            if current_chapter and chapter_content:
                chapters.append({
                    "title": current_chapter,
                    "content": "\n".join(chapter_content),
                    "knowledge_points": []
                })

            # 如果没有检测到章节，创建一个默认章节
            if not chapters:
                chapters.append({
                    "title": "课程内容",
                    "content": content[:2000],  # 限制长度
                    "knowledge_points": knowledge_points
                })

            # 为每章分配相关的知识点
            if knowledge_points:
                for i, chapter in enumerate(chapters):
                    # 简单的知识点分配：轮询分配
                    points_per_chapter = max(1, len(knowledge_points) // len(chapters))
                    start_idx = i * points_per_chapter
                    end_idx = start_idx + points_per_chapter
                    chapter["knowledge_points"] = knowledge_points[start_idx:end_idx]

            logger.info(f"内容分割完成，共{len(chapters)}章")
            return chapters

        except Exception as e:
            logger.error(f"章节分割失败：{str(e)}")
            # 返回单个章节
            return [{
                "title": "课程内容",
                "content": content[:2000],
                "knowledge_points": knowledge_points
            }]

    async def _save_parsed_content(self, material_id: str, structured_content: Dict[str, Any]) -> bool:
        """保存解析后的内容到数据库"""
        try:
            db = SessionLocal()

            # 删除旧的解析内容
            db.query(MaterialContent).filter(
                MaterialContent.material_id == material_id
            ).delete()

            # 保存新的章节内容
            chapters = structured_content.get("chapters", [])
            for idx, chapter in enumerate(chapters, 1):
                content = MaterialContent(
                    material_id=material_id,
                    chapter_number=idx,
                    chapter_title=chapter.get("title", ""),
                    content_text=chapter.get("content", "")[:5000],  # 限制长度
                    summary=structured_content.get("summary", ""),
                    knowledge_points=chapter.get("knowledge_points", [])
                )
                db.add(content)

            # 更新课件状态
            material = db.query(CourseMaterial).filter(
                CourseMaterial.id == material_id
            ).first()

            if material:
                material.upload_status = "completed"
                material.parsed_at = datetime.now()

            db.commit()
            db.close()

            logger.info(f"课件{material_id}的解析内容已保存到数据库")
            return True

        except Exception as e:
            logger.error(f"保存解析内容失败：{str(e)}")
            return False


# 全局解析器实例
_parser: Optional[MaterialParser] = None


def get_material_parser() -> MaterialParser:
    """获取课件解析器实例"""
    global _parser
    if _parser is None:
        _parser = MaterialParser()
    return _parser


async def parse_material_file(material_id: str, file_path: str) -> Dict[str, Any]:
    """解析课件文件的便捷函数"""
    parser = get_material_parser()
    return await parser.parse_material(material_id, file_path)